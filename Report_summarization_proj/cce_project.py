#"D:/ML-AI-iisc/Capstone_project/WIP/input_data" 
##import os
##import logging
##from transformers import AutoTokenizer, AutoModelForCausalLM, BitsAndBytesConfig
##from langchain_community.vectorstores import FAISS
##from langchain_community.embeddings import HuggingFaceEmbeddings
##from langchain.text_splitter import RecursiveCharacterTextSplitter
##from langchain.chains import RetrievalQA
##from langchain_community.llms import HuggingFacePipeline
##from pptx import Presentation
##from pptx.util import Inches
##import PyPDF2
##import fitz  # PyMuPDF for image extraction
##from accelerate import init_empty_weights, infer_auto_device_map
##from accelerate import infer_auto_device_map
##from transformers import AutoTokenizer, AutoModelForCausalLM, BitsAndBytesConfig
### **Logging Setup**
##logging.basicConfig(level=logging.INFO)
##logger = logging.getLogger(__name__)
##
### **Model Setup**
##model_name = "tiiuae/falcon-7b-instruct"
##
##logger.info("Starting to load tokenizer...")
##tokenizer = AutoTokenizer.from_pretrained(model_name)
##
### Initialize model
##print("Loading model...")
##bnb_config = BitsAndBytesConfig(
##    load_in_8bit=True  # Reduces memory usage
##)
##
### Load model without `device_map` initially
##model = AutoModelForCausalLM.from_pretrained(
##    model_name,
##    quantization_config=bnb_config
##)
##
### Infer device mapping
##device_map = infer_auto_device_map(model)
##
### Reload model with device mapping and offloading
##print("Reloading model with device map and offload settings...")
##model = AutoModelForCausalLM.from_pretrained(
##    model_name,
##    device_map=device_map,
##    offload_folder="./model_offload",
##    offload_state_dict=True,
##    quantization_config=bnb_config
##)
##
##print("Model loaded successfully!")
##
### **Embedding Setup**
##logger.info("Setting up embeddings and vector store...")
##embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
##vector_store = FAISS.load_local("./vector_store", embedding_model) if os.path.exists("./vector_store") else None
##logger.info("Vector store setup complete.")
##
### **PDF Extraction Function**
##def extract_pdf_content(folder_path):
##    """Extract text and images from PDFs in a folder."""
##    logger.info("Extracting content from PDFs...")
##    content = []
##    for file in os.listdir(folder_path):
##        if file.endswith(".pdf"):
##            file_path = os.path.join(folder_path, file)
##            logger.info(f"Processing {file}...")
##            
##            # Extract text
##            pdf_text = ""
##            with open(file_path, "rb") as f:
##                reader = PyPDF2.PdfReader(f)
##                for page in reader.pages:
##                    pdf_text += page.extract_text()
##            
##            # Extract images
##            images = []
##            pdf_document = fitz.open(file_path)
##            for page_num in range(len(pdf_document)):
##                page = pdf_document[page_num]
##                for img_index, img in enumerate(page.get_images(full=True)):
##                    xref = img[0]
##                    base_image = pdf_document.extract_image(xref)
##                    image_bytes = base_image["image"]
##                    image_path = os.path.join(folder_path, f"extracted_image_{file}_{page_num}_{img_index}.png")
##                    with open(image_path, "wb") as img_file:
##                        img_file.write(image_bytes)
##                    images.append(image_path)
##            
##            content.append({"text": pdf_text, "images": images})
##    logger.info("PDF content extraction complete.")
##    return content
##
### **Summarization Function**
##def summarize_content(content, llm):
##    """Summarize the extracted content using a language model."""
##    logger.info("Summarizing content...")
##    summaries = []
##    for doc in content:
##        input_text = doc["text"]
##        query = "Summarize this CAE report and identify failed components or unmet targets."
##        chain = RetrievalQA.from_chain_type(
##            llm=llm,
##            retriever=vector_store.as_retriever(),
##            return_source_documents=False
##        )
##        summary = chain.run(input_text)
##        summaries.append({"summary": summary, "images": doc["images"]})
##    logger.info("Summarization complete.")
##    return summaries
##
### **Create PPT Function**
##def create_ppt(summaries, output_path="summary_presentation.pptx"):
##    """Generate a PowerPoint presentation from summaries and images."""
##    logger.info("Creating PowerPoint presentation...")
##    ppt = Presentation()
##    for idx, summary in enumerate(summaries):
##        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
##        title = slide.shapes.title
##        title.text = f"Component Analysis {idx + 1}"
##        
##        # Add Summary
##        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
##        text_frame = text_box.text_frame
##        text_frame.text = summary["summary"]
##        
##        # Add Images
##        for img_idx, img_path in enumerate(summary["images"][:3]):  # Limit to 3 images per slide
##            left = Inches(1 + (img_idx * 2))
##            top = Inches(4)
##            slide.shapes.add_picture(img_path, left, top, width=Inches(2), height=Inches(2))
##    
##    ppt.save(output_path)
##    logger.info(f"Presentation saved to {output_path}")
##
### **Main Execution**
##if __name__ == "__main__":
##    pdf_folder = "D:/ML-AI-iisc/Capstone_project/WIP/input_data"  # Adjust path to folder with PDFs
##    output_ppt = "summary_presentation.pptx"
##    
##    # Extract PDF content
##    extracted_content = extract_pdf_content(pdf_folder)
##    
##    # Initialize LLM
##    pipeline = HuggingFacePipeline(model=model, tokenizer=tokenizer)
##    
##    # Summarize content
##    summarized_content = summarize_content(extracted_content, pipeline)
##    
##    # Create PPT
##    create_ppt(summarized_content, output_ppt)
##    logger.info("Process completed successfully.")
# working sumarization code
##import os
##from PyPDF2 import PdfReader
##from transformers import pipeline
##
### Directory containing PDF files
##pdf_folder = "D:/ML-AI-iisc/Capstone_project/WIP/input_data"
##
### Initialize the summarization pipeline
##summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
##
### Function to extract text from a PDF file
##def extract_text_from_pdf(pdf_path):
##    try:
##        reader = PdfReader(pdf_path)
##        text = ""
##        for page in reader.pages:
##            text += page.extract_text()
##        return text
##    except Exception as e:
##        print(f"Error reading {pdf_path}: {e}")
##        return None
##
### Function to summarize text
##def summarize_text(text, max_length=130, min_length=30):
##    try:
##        summary = summarizer(text, max_length=max_length, min_length=min_length, do_sample=False)
##        return summary[0]['summary_text']
##    except Exception as e:
##        print(f"Error summarizing text: {e}")
##        return None
##
### Process all PDF files in the folder
##for file_name in os.listdir(pdf_folder):
##    if file_name.endswith(".pdf"):
##        pdf_path = os.path.join(pdf_folder, file_name)
##        print(f"\nProcessing: {file_name}")
##
##        # Extract text from PDF
##        pdf_text = extract_text_from_pdf(pdf_path)
##        if not pdf_text:
##            print("Failed to extract text.")
##            continue
##
##        # Summarize text
##        summary = summarize_text(pdf_text)
##        if summary:
##            print("Summary:")
##            print(f"- {summary}")
##        else:
##            print("Failed to generate summary.")
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
from transformers import pipeline
from PIL import Image
from pdf2image import convert_from_path

# Directory containing PDF files
pdf_folder = "D:/ML-AI-iisc/Capstone_project/WIP/input_data"
ppt_output_path = "output_presentation.pptx"

# Initialize the summarization pipeline
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# Function to extract text from a PDF file
def extract_text_from_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        print(f"Error reading {pdf_path}: {e}")
        return None

# Function to summarize text
def summarize_text(text, max_length=130, min_length=30):
    try:
        summary = summarizer(text, max_length=max_length, min_length=min_length, do_sample=False)
        return summary[0]['summary_text']
    except Exception as e:
        print(f"Error summarizing text: {e}")
        return None

# Function to extract images from a PDF file
def extract_images_from_pdf(pdf_path, image_folder="temp_images"):
    os.makedirs(image_folder, exist_ok=True)
    try:
        images = convert_from_path(pdf_path)
        image_paths = []
        for i, image in enumerate(images):
            image_path = os.path.join(image_folder, f"page_{i+1}.png")
            image.save(image_path, "PNG")
            image_paths.append(image_path)
        return image_paths
    except Exception as e:
        print(f"Error extracting images from {pdf_path}: {e}")
        return []

# Function to create a PowerPoint presentation
def create_ppt(pdf_file, summary, image_paths, presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
    title = slide.shapes.title
    title.text = f"Summary for {os.path.basename(pdf_file)}"
    
    # Add summary bullet points
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(4.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    for point in summary.split(". "):  # Each sentence becomes a bullet point
        if point.strip():
            text_frame.add_paragraph().text = f"- {point.strip()}"

    # Add images
    for idx, image_path in enumerate(image_paths):
        if idx < 3:  # Limit to 3 images per slide to avoid overcrowding
            left = Inches(0.5 + (3 * idx))
            top = Inches(5.5)
            slide.shapes.add_picture(image_path, left, top, height=Inches(1.5))

# Main processing
presentation = Presentation()

for file_name in os.listdir(pdf_folder):
    if file_name.endswith(".pdf"):
        pdf_path = os.path.join(pdf_folder, file_name)
        print(f"Processing: {file_name}")
        
        # Extract text
        pdf_text = extract_text_from_pdf(pdf_path)
        if not pdf_text:
            print("Failed to extract text.")
            continue

        # Summarize text
        summary = summarize_text(pdf_text)
        if not summary:
            print("Failed to summarize text.")
            continue

        # Extract images
        images = extract_images_from_pdf(pdf_path)
        
        # Create PPT slide
        create_ppt(pdf_path, summary, images, presentation)

# Save the presentation
presentation.save(ppt_output_path)
print(f"Presentation saved to {ppt_output_path}")
